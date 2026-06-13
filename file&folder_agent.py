import argparse
import json
import os
import re
import sys
import textwrap
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import requests
from dotenv import load_dotenv

load_dotenv()

try:
    import docx as docx_lib
    HAVE_DOCX = True
except ImportError:
    HAVE_DOCX = False

try:
    from pptx import Presentation
    HAVE_PPTX = True
except ImportError:
    HAVE_PPTX = False

try:
    import openpyxl
    HAVE_XLSX = True
except ImportError:
    HAVE_XLSX = False

try:
    import pdfplumber
    HAVE_PDF = True
except ImportError:
    HAVE_PDF = False

OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
DEFAULT_MODEL   = os.getenv("OLLAMA_DEFAULT_MODEL", "llama3.2:latest")
MAX_CHARS       = 8000

TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm",
    ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".sh",
    ".yaml", ".yml", ".log", ".ini", ".cfg",
}

def ollama_chat(messages: list[dict], model: str) -> str:
    payload = {"model": model, "messages": messages, "stream": False}
    try:
        resp = requests.post(f"{OLLAMA_BASE_URL}/api/chat", json=payload, timeout=180)
        resp.raise_for_status()
        return resp.json()["message"]["content"]
    except requests.exceptions.ConnectionError:
        sys.exit(
            "[ERROR] Cannot reach Ollama at localhost:11434. "
            "Make sure 'ollama serve' is running."
        )

def _clean_json(raw: str) -> str:
    return re.sub(r"```json|```", "", raw).strip()

def read_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="replace")[:MAX_CHARS]
    except Exception as exc:
        return f"[ERROR reading text file: {exc}]"

def read_docx_file(path: Path) -> str:
    if not HAVE_DOCX:
        return "[python-docx not installed - run: pip install python-docx]"
    try:
        doc = docx_lib.Document(str(path))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                parts.append(" | ".join(cells))
        return "\n".join(parts)[:MAX_CHARS]
    except Exception as exc:
        return f"[ERROR reading docx: {exc}]"

def read_pptx_file(path: Path) -> str:
    if not HAVE_PPTX:
        return "[python-pptx not installed - run: pip install python-pptx]"
    try:
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs)
                        if text.strip():
                            slide_text.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        slide_text.append(" | ".join(cells))
            if slide_text:
                parts.append(f"[Slide {i}] " + " / ".join(slide_text))
        return "\n".join(parts)[:MAX_CHARS]
    except Exception as exc:
        return f"[ERROR reading pptx: {exc}]"

def read_xlsx_file(path: Path) -> str:
    if not HAVE_XLSX:
        return "[openpyxl not installed - run: pip install openpyxl]"
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            parts.append(f"[Sheet: {sheet.title}]")
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                if row_count >= 50:
                    parts.append("... (truncated)")
                    break
                if any(c is not None for c in row):
                    parts.append(" | ".join("" if c is None else str(c) for c in row))
                    row_count += 1
        return "\n".join(parts)[:MAX_CHARS]
    except Exception as exc:
        return f"[ERROR reading xlsx: {exc}]"

def read_csv_file(path: Path) -> str:
    return read_text_file(path)

def read_pdf_file(path: Path) -> str:
    if not HAVE_PDF:
        return "[pdfplumber not installed - run: pip install pdfplumber]"
    try:
        parts = []
        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(f"[Page {i}] {text}")
                if sum(len(p) for p in parts) > MAX_CHARS:
                    break
        return "\n".join(parts)[:MAX_CHARS]
    except Exception as exc:
        return f"[ERROR reading pdf: {exc}]"

READERS = {
    ".docx": read_docx_file,
    ".pptx": read_pptx_file,
    ".xlsx": read_xlsx_file,
    ".xlsm": read_xlsx_file,
    ".csv":  read_csv_file,
    ".pdf":  read_pdf_file,
}

def list_directory(path: str) -> dict:
    p = Path(path).expanduser().resolve()
    if not p.exists():
        return {"path": str(p), "exists": False, "is_dir": False, "folders": [], "files": []}
    if not p.is_dir():
        return {"path": str(p), "exists": True, "is_dir": False, "folders": [], "files": []}

    folders, files = [], []
    try:
        for entry in sorted(p.iterdir()):
            if entry.is_dir():
                folders.append(entry.name)
            else:
                stat = entry.stat()
                files.append({
                    "name": entry.name,
                    "ext": entry.suffix.lower(),
                    "size_bytes": stat.st_size,
                    "modified": datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc).isoformat(),
                })
    except PermissionError:
        pass
    return {"path": str(p), "exists": True, "is_dir": True, "folders": folders, "files": files}

def read_file_content(path: str) -> dict:
    p = Path(path).expanduser().resolve()
    if not p.exists() or not p.is_file():
        return {"path": str(p), "exists": False, "ext": "", "size_bytes": 0, "content": "", "truncated": False}

    ext = p.suffix.lower()
    size = p.stat().st_size
    if ext in READERS:
        content = READERS[ext](p)
    elif ext in TEXT_EXTENSIONS:
        content = read_text_file(p)
    else:
        content = f"[Unsupported file type: {ext or '(no extension)'} - skipped]"

    return {
        "path": str(p),
        "exists": True,
        "ext": ext,
        "size_bytes": size,
        "content": content,
        "truncated": len(content) >= MAX_CHARS,
    }

def walk_folder(path: str, max_files: int = 50, max_depth: int = 5) -> list[dict]:
    p = Path(path).expanduser().resolve()
    results = []
    if not p.exists() or not p.is_dir():
        return results

    base_depth = len(p.parts)
    for root, dirs, files in os.walk(p):
        depth = len(Path(root).parts) - base_depth
        if depth >= max_depth:
            dirs[:] = []
            continue

        dirs[:] = [d for d in dirs if not d.startswith(".")]
        for fname in sorted(files):
            if fname.startswith("."):
                continue
            if len(results) >= max_files:
                return results

            file_path = Path(root) / fname
            ext = file_path.suffix.lower()
            if ext in READERS or ext in TEXT_EXTENSIONS:
                results.append(read_file_content(str(file_path)))

    return results

def parse_intent(user_command: str, model: str) -> dict:
    system_prompt = textwrap.dedent("""
        You are an intent-extraction assistant for a local file-browsing agent.

        Given the user command, output ONLY a single JSON object (no markdown, no prose)
        with these keys:

        action (string, required):
            "list_dir"       - just list the contents of a folder (no reading file content)
            "read_file"      - read and analyse ONE specific file
            "analyze_folder" - recursively read files in a folder and analyse them together

        path (string, required): the filesystem path mentioned by the user.
            If relative or ambiguous, use it as given (e.g. "./reports", "C:\\Users\\me\\Desktop").

        task (string, required):
            "summarize" - summarise the content
            "analyze"   - deeper analysis (themes, structure, insights)
            "list"      - just enumerate contents, no content analysis needed
            "find_info" - search for specific information the user asked about

        Output ONLY the JSON object.
    """).strip()

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user",   "content": user_command},
    ]
    raw = ollama_chat(messages, model)
    raw = _clean_json(raw)

    try:
        return json.loads(raw)
    except json.JSONDecodeError as exc:
        raise ValueError(f"Ollama returned invalid JSON for intent:\n{raw}") from exc

def _file_blob(files: list[dict]) -> str:
    parts = []
    for f in files:
        if not f.get("exists"):
            continue
        parts.append(
            f"--- FILE: {f['path']} (ext={f['ext']}, size={f['size_bytes']} bytes) ---\n"
            f"{f['content'][:3000]}"
        )
    return "\n\n".join(parts)

def analyze_listing(listing: dict, task: str, model: str) -> dict:
    system_prompt = (
        "You are a file-system assistant. You receive a directory listing "
        "(folder names and file metadata) and produce a structured JSON summary."
    )
    user_prompt = textwrap.dedent(f"""
        Directory: {listing['path']}
        Folders: {listing['folders']}
        Files: {json.dumps(listing['files'], indent=2)}

        Task: {task}

        Return ONLY a JSON object with keys:
          "summary"       - 1-3 sentence overview of what's in this folder
          "folder_count"  - number of subfolders
          "file_count"    - number of files
          "file_types"    - object mapping extension -> count
          "notable_items" - list of up to 5 notable file/folder names with a short reason
    """).strip()

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user",   "content": user_prompt},
    ]
    raw = _clean_json(ollama_chat(messages, model))
    try:
        return json.loads(raw)
    except json.JSONDecodeError as exc:
        raise ValueError(f"Ollama returned invalid JSON for listing analysis:\n{raw}") from exc

def analyze_files(files: list[dict], task: str, model: str, user_command: str) -> dict:
    system_prompt = (
        "You are a document analysis assistant. You receive the extracted text "
        "content of one or more files and the user's original request. "
        "Return ONLY valid JSON, no markdown, no prose outside the JSON."
    )

    if task == "summarize":
        instructions = (
            'Return a JSON object: {"summary": "...", "files": '
            '[{"path": "...", "summary": "..."}]}'
        )
    elif task == "find_info":
        instructions = (
            'Find information relevant to the user\'s request. Return a JSON object: '
            '{"answer": "...", "evidence": [{"path": "...", "excerpt": "...", "relevance": "..."}]}'
        )
    elif task == "analyze":
        instructions = (
            'Provide deeper analysis. Return a JSON object: '
            '{"overview": "...", "themes": ["..."], "files": '
            '[{"path": "...", "summary": "...", "key_points": ["..."]}], "insights": ["..."]}'
        )
    else:
        instructions = (
            'Return a JSON object: {"files": [{"path": "...", "title_or_topic": "..."}]}'
        )

    user_prompt = (
        f"User request: {user_command}\n\n"
        f"Task: {task}\n\n"
        f"Instructions: {instructions}\n\n"
        f"Files:\n\n{_file_blob(files)}"
    )

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user",   "content": user_prompt},
    ]
    raw = _clean_json(ollama_chat(messages, model))
    try:
        return json.loads(raw)
    except json.JSONDecodeError as exc:
        raise ValueError(f"Ollama returned invalid JSON for file analysis:\n{raw}") from exc

def run_agent(user_command: str, model: str) -> dict[str, Any]:
    print(f"[Agent] Parsing intent for: {user_command!r}")
    intent = parse_intent(user_command, model)
    print(f"[Agent] Intent -> {json.dumps(intent, indent=2)}")

    action = intent.get("action", "list_dir")
    path   = intent.get("path", ".")
    task   = intent.get("task", "summarize")

    result: dict[str, Any] = {
        "command":   user_command,
        "intent":    intent,
        "timestamp": datetime.now(timezone.utc).isoformat(),
    }

    if action == "list_dir":
        listing = list_directory(path)
        result["listing"] = listing
        if listing["exists"] and listing["is_dir"]:
            print("[Agent] Analysing directory listing ...")
            result["analysis"] = analyze_listing(listing, task, model)
        else:
            result["analysis"] = {"error": f"Path not found or not a directory: {path}"}
    elif action == "read_file":
        file_data = read_file_content(path)
        result["file"] = {k: v for k, v in file_data.items() if k != "content"}
        if file_data["exists"]:
            print("[Agent] Analysing file content ...")
            result["analysis"] = analyze_files([file_data], task, model, user_command)
        else:
            result["analysis"] = {"error": f"File not found: {path}"}
    elif action == "analyze_folder":
        files = walk_folder(path)
        result["files_found"] = len(files)
        result["files"] = [{k: v for k, v in f.items() if k != "content"} for f in files]
        if files:
            print(f"[Agent] Analysing {len(files)} file(s) from folder ...")
            result["analysis"] = analyze_files(files, task, model, user_command)
        else:
            result["analysis"] = {"error": f"No readable files found in: {path}"}
    else:
        result["analysis"] = {"error": f"Unknown action: {action}"}

    return result

def main():
    parser = argparse.ArgumentParser(description="Local file/folder browsing & analysis agent")
    parser.add_argument("--model",  default=DEFAULT_MODEL, help=f"Ollama model name (default: {DEFAULT_MODEL})")
    parser.add_argument("--output", default=None, help="Write JSON result to this file")
    args = parser.parse_args()

    print("File Agent  |  read-only, browses files/folders on this machine")
    print("Examples:")
    print("  list the files in ./documents")
    print("  summarize report.pdf")
    print("  analyze the folder ./project_docs and tell me the main themes")
    print("type 'exit' or Ctrl-C to quit\n")

    while True:
        try:
            command = input("You: ").strip()
        except (KeyboardInterrupt, EOFError):
            print("\nBye!")
            break

        if not command:
            continue
        if command.lower() in {"exit", "quit", "q"}:
            print("Bye!")
            break

        try:
            result = run_agent(command, model=args.model)
            pretty = json.dumps(result, indent=2, ensure_ascii=True)
            print("\n-- Agent Response -------------------------------------")
            print(pretty)
            print("-------------------------------------------------------\n")

            if args.output:
                with open(args.output, "w") as fh:
                    fh.write(pretty)
                print(f"[Agent] Result written to {args.output}\n")

        except Exception as exc:
            clean_exc = str(exc).encode('ascii', 'ignore').decode('ascii')
            print(f"[ERROR] {clean_exc}\n")

if __name__ == "__main__":
    main()